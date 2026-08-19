# 발표자료(PPTX) 생성 — 과제 제출용 슬라이드.
#
# 기존 과제제출보고서.pdf(문서형, 48쪽)와는 별도로, 발표 형식(16:9 슬라이드)으로
# 프로젝트를 요약한다. 게임의 실제 색상 토큰(다크 배경 + 민트 브랜드)을 그대로 써서
# 스크린샷과 슬라이드가 한 톤으로 보이게 한다.
#
# 실행: pip install python-pptx && python tools/build_slides.py
# PDF 변환: python tools/pptx_to_pdf.py  (PowerPoint COM 자동화, 별도 스크립트)

import os, sys, csv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

sys.stdout.reconfigure(encoding='utf-8')
os.chdir(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ── 색상 토큰 (게임 CSS와 동일) ─────────────────────────────
BG      = RGBColor(0x0E, 0x11, 0x16)
BG_SOFT = RGBColor(0x16, 0x1B, 0x22)
LINE    = RGBColor(0x23, 0x2A, 0x34)
INK     = RGBColor(0xE8, 0xED, 0xF4)
INK_DIM = RGBColor(0x8A, 0x94, 0xA6)
BRAND   = RGBColor(0x5E, 0xE6, 0xC5)
WARN    = RGBColor(0xFF, 0xB4, 0x54)
BAD     = RGBColor(0xFF, 0x6B, 0x6B)

URL = "https://odd-one-out-nine.vercel.app/"
REPO = "https://github.com/stulss/odd-one-out"

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf

def para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return p

def bullets(s, l, t, w, h, items, size=18, color=INK, gap=10):
    tb, tf = box(s, l, t, w, h)
    for i, it in enumerate(items):
        lvl = 0
        txt = it
        if isinstance(it, tuple): txt, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("▸ " if lvl == 0 else "· ") + txt
        r.font.size = Pt(size - lvl * 2); r.font.color.rgb = color if lvl == 0 else INK_DIM
        r.font.name = FONT
    return tb

def title_bar(s, kicker, title, color=BRAND):
    _, tf = box(s, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.4))
    para(tf, kicker, 14, color=color, bold=True, first=True)
    _, tf2 = box(s, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.9))
    para(tf2, title, 32, color=INK, bold=True, first=True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.55), Inches(11.9), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False

def page_no(s, n):
    _, tf = box(s, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
    para(tf, str(n), 11, color=INK_DIM, align=PP_ALIGN.RIGHT, first=True)

def pill(s, l, t, w, h, text, fill=BG_SOFT, color=INK, size=13, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return sh

def img(s, path, l, t, h):
    if not os.path.exists(path): return None
    from PIL import Image
    im = Image.open(path); ar = im.width / im.height
    return s.shapes.add_picture(path, l, t, height=h, width=Emu(int(h * ar)))

def phone_row(s, items, top, h=Inches(4.1), gap=Inches(0.35), start_l=None, max_right=None):
    """items: [(경로, 캡션), ...] — 세로형 스크린샷을 가로로 나열.
    지정한 높이로 쟀을 때 오른쪽 여백(max_right, 기본 슬라이드 폭-0.7")을 넘으면
    전체 행이 폭에 맞게 자동으로 줄어든다 (수동으로 h를 맞출 필요가 없다)."""
    n = len(items)
    from PIL import Image
    ars = []
    for p, _ in items:
        if os.path.exists(p):
            im = Image.open(p); ars.append(im.width / im.height)
        else:
            ars.append(0.462)
    l0 = start_l if start_l is not None else Inches(0.7)
    right_limit = max_right if max_right is not None else (W - Inches(0.7))
    avail = right_limit - l0 - gap * (n - 1)
    max_h_by_width = Emu(int(avail / sum(ars))) if sum(ars) > 0 else h
    use_h = min(h, max_h_by_width)
    widths = [int(use_h * ar) for ar in ars]
    total = sum(widths) + gap * (n - 1)
    l = start_l if start_l is not None else int((W - total) / 2)
    for (p, cap), wpx in zip(items, widths):
        img(s, p, l, top, use_h)
        _, tf = box(s, l - Inches(0.2), top + use_h + Inches(0.06), Emu(wpx) + Inches(0.4), Inches(0.35))
        para(tf, cap, 12, color=INK_DIM, align=PP_ALIGN.CENTER, first=True)
        l += wpx + gap
    return use_h

def table_slide(s, headers, rows, l, t, w, h, col_w=None, header_fill=LINE, font_size=12, left_cols=()):
    """left_cols: 왼쪽 정렬할 열 인덱스 집합. 문장형 텍스트(항목·증빙자료 등)는
    가운데 정렬하면 줄마다 들쭉날쭉해 보여서 왼쪽 정렬이 낫다."""
    rc, cc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(rc, cc, l, t, w, h).table
    if col_w:
        for i, cw in enumerate(col_w): gt.columns[i].width = cw
    for c, htext in enumerate(headers):
        cell = gt.cell(0, c); cell.text = str(htext)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        align = PP_ALIGN.LEFT if c in left_cols else PP_ALIGN.CENTER
        for p in cell.text_frame.paragraphs:
            p.alignment = align
            for r in p.runs: r.font.bold = True; r.font.size = Pt(font_size); r.font.color.rgb = INK; r.font.name = FONT
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(ri, c); cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = BG_SOFT
            cell.margin_left = Pt(8) if c in left_cols else cell.margin_left
            align = PP_ALIGN.LEFT if c in left_cols else PP_ALIGN.CENTER
            for p in cell.text_frame.paragraphs:
                p.alignment = align
                for r in p.runs: r.font.size = Pt(font_size); r.font.color.rgb = INK_DIM; r.font.name = FONT
    return gt

def evidence_cell(s, x, y, w, h, num, title, imgs, note=None):
    """체크리스트 항목 1개를 '번호+제목 / 스크린샷 1~2장 / 캡션' 카드로 그린다.
    imgs: [(경로, 캡션), ...] 최대 2장. 세로형·가로형 스크린샷 모두 지원."""
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = BG_SOFT
    card.line.color.rgb = LINE; card.line.width = Pt(0.75); card.shadow.inherit = False

    _, tf = box(s, x + Inches(0.12), y + Inches(0.06), w - Inches(0.24), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"{num}. {title}"
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = BRAND; r.font.name = FONT

    if imgs:
        from PIL import Image
        img_top = y + Inches(0.5)
        img_h = h - Inches(0.5) - (Inches(0.32) if note else Inches(0.12))
        ars = []
        for p_, _ in imgs:
            if os.path.exists(p_):
                im = Image.open(p_); ars.append(im.width / im.height)
            else:
                ars.append(0.462)
        gap = Inches(0.12)
        avail = w - Inches(0.24) - gap * (len(imgs) - 1)
        max_h_by_w = Emu(int(avail / sum(ars))) if sum(ars) > 0 else img_h
        use_h = min(img_h, max_h_by_w)
        widths = [int(use_h * ar) for ar in ars]
        total = sum(widths) + gap * (len(imgs) - 1)
        l = x + (w - total) / 2
        for (p_, cap), wpx in zip(imgs, widths):
            img(s, p_, l, img_top, use_h)
            l += wpx + gap
        note_top = img_top + use_h + Inches(0.03)
    else:
        note_top = y + Inches(0.55)

    if note:
        _, tf2 = box(s, x + Inches(0.12), note_top, w - Inches(0.24), Inches(0.3))
        para(tf2, note, 10, color=INK_DIM, align=PP_ALIGN.CENTER, first=True, space_after=0)

def read_csv(p):
    if not os.path.exists(p): return []
    with open(p, encoding='utf-8-sig') as f:
        return list(csv.reader(f))

N = [0]
def track(s):
    N[0] += 1
    page_no(s, N[0])
    return s

# ═══════════════════════════════════════════════════════════
# 1. 표지
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
# 미니 격자 아이콘 (타이틀 화면 UI를 축소 재현)
gx, gy, cell, gap = Inches(5.7), Inches(1.3), Inches(0.55), Inches(0.12)
colors4 = [RGBColor(0x3D,0x6B,0xFF)]*3 + [BRAND]
for i in range(4):
    r, c = divmod(i, 2)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        gx + c*(cell+gap), gy + r*(cell+gap), cell, cell)
    sh.fill.solid(); sh.fill.fore_color.rgb = colors4[i]; sh.line.fill.background(); sh.shadow.inherit = False
    if i == 3:
        sh.rotation = 12
_, tf = box(s, Inches(1.5), Inches(3.15), Inches(10.3), Inches(1.3))
para(tf, "딱 하나 이상함", 54, color=INK, bold=True, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, Inches(1.5), Inches(4.15), Inches(10.3), Inches(0.6))
para(tf, "여러 개 중 딱 하나 다른 것을 찾는 브라우저 관찰력 게임", 20, color=INK_DIM, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, Inches(1.5), Inches(4.95), Inches(10.3), Inches(0.5))
para(tf, "「내가 설계한 미니게임」 과제 제출 · 발표자료", 15, color=BRAND, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7))
para(tf, URL, 15, color=INK_DIM, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(tf, REPO, 12, color=INK_DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 2. 프로젝트 개요 + 완료기준
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "OVERVIEW", "프로젝트 개요")
bullets(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(4.5), [
    "격자 안 도형 여럿 중 딱 하나만 색·크기·회전 등이 미묘하게 다름",
    "제한 시간 안에 찾아 선택 → 다음 단계, 시간 초과 시 종료",
    ("차이 축 10종 · 도형 4종 · 격자 상한 4×6(24칸)", 1),
    "프레임워크·이미지·음원 파일 없이 순수 HTML/CSS/JS",
    ("소스 전체 80KB · 라이브러리 0개", 1),
    "같은 링크 = 같은 문제 (시드 난수, 서버 없음)",
], size=19)
pill_y = Inches(1.9)
labels = ["규칙·조작·상태\n상시 표시", "조작 1회 =\n결과 1회", "다시 시작\n= 초기화", "음소거·움직임\n줄이기 작동"]
for i, lab in enumerate(labels):
    pill(s, Inches(7.0), pill_y + i*Inches(1.15), Inches(5.6), Inches(0.95), lab, fill=BG_SOFT, color=BRAND, size=15)
_, tf = box(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.35))
para(tf, "과제 완료기준 4항목 — 전부 충족", 14, color=INK_DIM, first=True)

# ═══════════════════════════════════════════════════════════
# 2b. 기능 안내서
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "USER GUIDE", "기능 안내서 — 플레이 방법과 전체 기능")
rows = [
    ["목표", "다른 하나 찾기", "격자 안 도형 여럿 중 색·크기·회전 등이 미묘하게 다른 딱 하나를 제한 시간 안에 선택"],
    ["조작", "클릭·탭·방향키+Enter · P(일시정지)", "마우스·터치·키보드만으로 완주 가능, 화면 이탈(blur) 시 자동 일시정지"],
    ["진행", "정답 → 다음 단계 / 오답 → 시간 1초 차감", "오답은 즉시 게임오버가 아님 — 흔들림 효과 후 계속 진행"],
    ["오늘의 문제", "/?d=20260818", "KST 자정 기준 하루 1문제 — 전 세계 누구나 같은 문제"],
    ["도전장", "/?c=K3F9QZ", "친구와 완전히 동일한 문제를 공유, 상대 기록(&s=12) 배너 표시"],
    ["설정", "소리·움직임 줄이기·진동·색약 모드·난이도(T)", "클릭 즉시 반영, 재시작 불필요, 새로고침 후에도 유지"],
    ["기록", "최고 단계·점수·연속 출석", "localStorage 저장, 값이 손상돼도 자동으로 기본값 복구"],
    ["결과 공유", "결과 공유 버튼 1회 클릭", "네이티브 공유 → 클립보드 이미지 복사 → 폴백 3버튼 순으로 시도"],
]
table_slide(s, ["구분", "기능 / 조작", "설명"], rows,
            Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.9),
            col_w=[Inches(1.7), Inches(3.4), Inches(6.8)], font_size=13, left_cols={1, 2})

# ═══════════════════════════════════════════════════════════
# 3. 핵심 규칙 (카드1) — 스크린샷
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 1 · 30초 규칙", "조작 1회 = 결과 1회")
bullets(s, Inches(0.7), Inches(1.9), Inches(4.6), Inches(4.5), [
    "플레이어는 격자 칸을 조작해 다른 하나를 선택하면 성공,",
    ("제한 시간이 0이 되면 실패", 1),
    "정답: 팝 애니메이션 + 점수 상승 + 다음 단계",
    "오답: 흔들림 + 남은 시간 1초 차감 (즉사 아님)",
    "화면 하단 3줄 상시 표시 — 규칙 · 조작 · 현재 상태",
], size=18)
phone_row(s, [
    ("img/actions/01_start_after.png", "게임 시작"),
    ("img/actions/02_correct_after.png", "정답 → 다음 단계"),
    ("img/actions/03_wrong_after.png", "오답 → 시간 차감"),
], top=Inches(1.85), h=Inches(4.6), start_l=Inches(6.9))

# ═══════════════════════════════════════════════════════════
# 4. 기술 선택
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "TECH DECISION", "게임 엔진을 쓰지 않은 이유")
_, tf = box(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.9))
para(tf, "“이 게임의 화면은 정지된 도형 격자 + 한 번의 입력이다.", 19, color=INK, first=True, space_after=2)
para(tf, "렌더 루프가 필요 없는 게임에 엔진을 쓰면 엔진의 비용만 남고 이득은 0이다.”", 19, color=INK)
rows = [["Vanilla + DOM", "약 25KB", "약 0.8초", "채택"],
        ["PixiJS", "약 140KB", "약 1.8초", "과함"],
        ["Phaser 3", "약 450KB", "약 3.0초", "목표 위반"],
        ["Unity WebGL", "15MB+", "약 15초", "치명적"]]
table_slide(s, ["후보", "최소 전송 크기", "첫 플레이까지", "판단"], rows,
            Inches(0.7), Inches(3.0), Inches(11.9), Inches(2.4),
            col_w=[Inches(3.5), Inches(2.8), Inches(2.8), Inches(2.8)])
bullets(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.3), [
    "이미지 파일 0개 — 미묘한 차이를 1px·3° 단위로 정확히 제어해야 난이도가 성립.",
    ("래스터 이미지는 축소 과정에서 차이가 뭉개지거나 노이즈가 생겨 불공정한 실패를 만든다.", 1),
], size=15)

# ═══════════════════════════════════════════════════════════
# 5. 아키텍처
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "ARCHITECTURE", "13개 모듈 구조")
mods = [
    ("game.js", "상태 머신"), ("rng.js", "시드 난수"), ("stage.js", "스테이지 생성"),
    ("render.js", "DOM 그리드"), ("input.js", "입력 처리"), ("timer.js", "rAF 루프 1개"),
    ("save.js", "손상 복구"), ("daily.js", "오늘의 문제"), ("card.js", "결과 카드"),
    ("share.js", "공유"), ("audio.js", "합성음"), ("logger.js", "플레이 기록"),
]
cols, rows_n = 4, 3
cw, ch, gx0, gy0, ggap = Inches(2.85), Inches(1.15), Inches(0.7), Inches(1.9), Inches(0.18)
for i, (name, desc) in enumerate(mods):
    r, c = divmod(i, cols)
    x = gx0 + c*(cw+ggap); y = gy0 + r*(ch+ggap)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_SOFT
    sh.line.color.rgb = LINE; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r_ = p.add_run(); r_.text = name
    r_.font.size = Pt(16); r_.font.bold = True; r_.font.color.rgb = BRAND; r_.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(12); r2.font.color.rgb = INK_DIM; r2.font.name = FONT
bullets(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0), [
    "설계 원칙 3가지 — rAF 루프 앱 전체 1개 · DOM 노드 최초 1회 생성 후 재사용 · 현재 판 상태는 저장하지 않음",
], size=15, color=INK_DIM)

# ═══════════════════════════════════════════════════════════
# 6~8. 문제 해결 Before/After 3건
# ═══════════════════════════════════════════════════════════
def problem_slide(kicker, title, before, after, note=None):
    s = track(new_slide())
    title_bar(s, kicker, title, color=WARN)
    _, tf = box(s, Inches(0.7), Inches(1.85), Inches(5.7), Inches(4.6))
    para(tf, "BEFORE", 14, color=BAD, bold=True, first=True, space_after=8)
    for i, line in enumerate(before):
        para(tf, "• " + line, 16, color=INK_DIM, space_after=8)
    _, tf2 = box(s, Inches(6.7), Inches(1.85), Inches(5.9), Inches(4.6))
    para(tf2, "AFTER", 14, color=BRAND, bold=True, first=True, space_after=8)
    for i, line in enumerate(after):
        para(tf2, "• " + line, 16, color=INK, space_after=8)
    if note:
        _, tf3 = box(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.7))
        para(tf3, note, 14, color=INK_DIM, first=True)
    return s

problem_slide("문제 발견 #1", "같은 링크인데 기기마다 다른 문제가 나왔다",
    ["격자 칸 수를 화면 크기에 맞춰 계산",
     "1280px → 6×8=48칸, 375px → 4×6=24칸",
     "정답 인덱스까지 달라짐 → '오늘의 문제' 성립 불가"],
    ["그리드를 스테이지 번호로만 결정 (화면 크기 무관)",
     "상한 4×6=24칸 고정 (360px에서도 44px 터치 보장)",
     "같은 시드 25스테이지 → 완전 일치 (실측)"],
    "교훈: 되돌리기 비싼 결함일수록 일찍 찾아야 한다 — 수정 규모는 코드 3줄이었다.")

problem_slide("문제 발견 #2", "통과하던 검사가 사실 아무것도 검사하지 않고 있었다",
    ["가로 넘침 검사: scrollWidth > innerWidth",
     "CSS에 overflow-x:hidden이 걸려 있어",
     "넘쳐도 항상 false → '검사 통과'가 거짓 신호"],
    ["게임판 좌표를 뷰포트 경계와 직접 비교하도록 변경",
     "390×844에서 11스테이지 전 구간 재측정",
     "게임판 16~374px, 실제로 넘침 없음 확인"],
    "교훈: 초록불을 믿기 전에 그 검사가 빨간불이 될 수 있는지 먼저 확인해야 한다.")

problem_slide("문제 발견 #3", "화면에 내놓은 난이도 손잡이가 거의 작동하지 않았다",
    ["난이도 T를 '제한 시간' 계수에만 연결",
     "자동 플레이 30판: T 3.6→5.4(50%↑)에도",
     "도달 단계 중앙값 21.5→20.0 (1.5단계 차이뿐)"],
    ["원인: 시간이 부족해지는 시점(30단계)이",
     "판이 끝나는 시점(20~25단계)보다 늦었다",
     "손잡이를 '차이 곡선'으로 이동 → 18~25%p 차이로 개선"],
    "교훈: 값을 만들었으면 흔들어보고 결과가 실제로 움직이는지 재봐야 한다.")

# ═══════════════════════════════════════════════════════════
# 9. 카드2 검사표
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 2 · PC 숨은 검사", "검사표 — 7/7 통과")
rows = [
    ["1366×768 (최악 격자 4×6)", "잘림·가로 넘침 없음", "넘침 0, 셀 90px"],
    ["1920×1080 (최악 격자 4×6)", "게임판 유지, 여백만 증가", "넘침 0, 셀 132px"],
    ["1366↔1920 창 크기 변경", "레이아웃 즉시 재계산", "리사이즈 즉시 반영"],
    ["1초에 10회 연속 클릭", "조작 1회 = 결과 1회", "정확히 1스테이지"],
    ["blur(포커스 이탈)", "자동 일시정지, 시간 보존", "즉시 전환·재개 확인"],
    ["P키 일시정지 → 재개", "타이머 정지·복원", "오차 없이 이어짐"],
    ["10분 연속 실행", "FPS 유지·메모리 증가 없음", "PASS (다음 슬라이드)"],
]
table_slide(s, ["검사 조건", "예상 결과", "실제 결과"], rows,
            Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.6),
            col_w=[Inches(4.3), Inches(4.0), Inches(3.6)], font_size=14)

# ═══════════════════════════════════════════════════════════
# 10. 10분 안정성
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 2 · 10분 실행", "10분 연속 자동 진행 — PASS")
phone_row(s, [
    ("img/stability/t000_start.png", "0분"),
    ("img/stability/t180_3min.png", "3분"),
    ("img/stability/t360_6min.png", "6분"),
    ("img/stability/t600_10min.png", "10분"),
], top=Inches(1.85), h=Inches(3.5))  # start_l 생략 → 슬라이드 폭 기준 가운데 정렬
metrics = [("DOM 노드", "112 → 112", "증가 0"), ("JS 힙", "9.54MB → 9.54MB", "증가 0"),
           ("프레임 드랍", "0 / 6002", "0.00%"), ("콘솔 오류", "0건", "10분 내내")]
mx = Inches(0.9)
for label, val, sub in metrics:
    pill(s, mx, Inches(5.75), Inches(2.75), Inches(1.05),
         f"{label}\n{val}\n{sub}", fill=BG_SOFT, color=BRAND, size=13)
    mx += Inches(2.9)

# ═══════════════════════════════════════════════════════════
# 11. 카드3 난이도 데이터
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 3 · 데이터 기반 조정", "난이도 T 전/후 비교")
_, tf = box(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.5))
para(tf, "⚠ 자동 플레이 데이터입니다 — 사람 반응·피로가 빠져 있음을 숨기지 않고 명시합니다.", 13,
     color=WARN, first=True)
phone_row(s, [
    ("img/difficulty/settings_T4_4.png", "T = 4.4 (기존값)"),
    ("img/difficulty/result_T4_4.png", "결과 (T=4.4)"),
    ("img/difficulty/settings_T5_4.png", "T = 5.4 (비교값)"),
    ("img/difficulty/result_T5_4.png", "결과 (T=5.4)"),
], top=Inches(2.35), h=Inches(3.15), start_l=Inches(0.9))
rows = [["도달 단계 중앙값", "21.0단계", "20.0단계"], ["생존 시간 중앙값", "27.9초", "25.4초"]]
table_slide(s, ["지표", "T=4.4 (전)", "T=5.4 (후)"], rows,
            Inches(9.6), Inches(2.4), Inches(3.0), Inches(1.3), font_size=12)
_, tf2 = box(s, Inches(9.6), Inches(3.9), Inches(3.0), Inches(2.6))
para(tf2, "최종값: T=4.4 유지", 14, color=BRAND, bold=True, first=True, space_after=6)
para(tf2, "단계 차이는 작지만 내부 지표(지각 여유)로는 20~25단계 구간에서 18~25%p 차이 확인.", 12, color=INK_DIM)

# ═══════════════════════════════════════════════════════════
# 12. 카드4 저장과 복구
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 4 · 저장과 복구", "손상값 6종 — 전부 기본값 복구")
cases = ["정상값", "빈 값", '"abc" (글자만)', "필수 항목 누락", "타입 오류", "배열"]
cx = Inches(0.7)
for c in cases:
    pill(s, cx, Inches(2.0), Inches(1.85), Inches(0.7), c, fill=BG_SOFT, color=INK, size=13)
    cx += Inches(1.98)
bullets(s, Inches(0.7), Inches(3.1), Inches(11.9), Inches(3.5), [
    "현재 판(점수·스테이지·타이머)은 저장하지 않는다 → 새 게임에서 자동 초기화",
    ("초기화 로직을 따로 만들면 반드시 빠뜨리는 항목이 생긴다는 판단", 1),
    "최고 기록·설정만 별도 키로 보존",
    "저장값이 6가지 방식으로 깨져도 콘솔 오류 0건, 항상 기본값으로 정상 실행",
], size=18)

# ═══════════════════════════════════════════════════════════
# 13. 카드5 손맛과 선택권
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "CARD 5 · 손맛과 선택권", "효과 + 접근성")
bullets(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(4.5), [
    "정답 선택 시 팝 애니메이션 + 효과음(WebAudio 합성, 파일 0개)",
    "음소거 → 클릭 즉시 무음 전환",
    "움직임 줄이기 → 흔들림·팝 대신 색 점멸로 대체",
    ("효과를 완전히 없애지 않는 이유: 사건 인지 자체는 유지해야 하기 때문", 1),
    "설정은 즉시 반영, 재시작 불필요",
], size=18)
phone_row(s, [
    ("img/actions/06_sound_after.png", "소리 끄기"),
    ("img/actions/07_motion_after.png", "움직임 줄이기"),
], top=Inches(1.85), h=Inches(4.6), start_l=Inches(7.6))

# ═══════════════════════════════════════════════════════════
# 14. 결과 공유
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "FEATURE", "결과 공유 — 클릭 한 번, 이미지 즉시 복사")
bullets(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(4.5), [
    "[결과 공유] 클릭 시 우선순위: 네이티브 공유 → 클립보드 이미지 복사 → 폴백",
    "데스크탑 대부분은 navigator.share가 없어 예전엔 매번 폴백 3버튼만 노출",
    "이제는 결과 카드 이미지를 클립보드에 바로 복사 → Ctrl+V로 즉시 붙여넣기",
    ("텍스트(링크)도 같은 클립보드 항목에 함께 담아 이미지 미지원 대상에도 대응", 1),
    "부수 발견: #fallback에 [hidden] CSS 누락 — 공유 전에도 폴백 버튼이 항상 보이던 버그 수정",
], size=17)
phone_row(s, [
    ("img/actions/09_share_after.png", "네이티브 공유 없음"),
    ("img/actions/10_copytext_after.png", "복사 안내 문구"),
], top=Inches(1.85), h=Inches(4.6), start_l=Inches(7.6))

# ═══════════════════════════════════════════════════════════
# 15. AI 협업 방식
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "AI 3줄", "맡긴 일 · 판단한 일 · 말을 안 들은 일")
labels = ["AI에게 맡긴 일", "내가 판단한 일", "AI 말을 안 들은 일"]
texts = [
    "코드 전체(13개 모듈) 작성과 검증 실행 — 시드 재현성·저장 손상 6종·연타 방지·해상도 등을 직접 측정",
    "배포처를 Cloudflare→GitHub Pages로 변경, 단계별 폴더 구성, 스크린샷 3상태 규칙, STEP 6 중단 판단",
    "step2 스크린샷 오류를 직접 열어보고 발견, 제출용 PDF에서 빠진 항목 5개를 찾아냄, 자동 플레이 데이터 오남용 방지",
]
ty = Inches(1.9)
for lab, txt in zip(labels, texts):
    pill(s, Inches(0.7), ty, Inches(3.0), Inches(0.5), lab, fill=LINE, color=BRAND, size=14)
    _, tf = box(s, Inches(3.9), ty - Inches(0.05), Inches(8.7), Inches(1.2))
    para(tf, txt, 15, color=INK_DIM, first=True)
    ty += Inches(1.55)

# ═══════════════════════════════════════════════════════════
# 16. 트러블슈팅 하이라이트
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "TROUBLESHOOTING", "실제로 겪은 문제 8건 중 대표 사례")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.6), [
    "같은 실수를 세 번 반복 — step2 재촬영 누락 · PDF 스크린샷 5개 누락 · 이모지 깨짐",
    ("세 번 다 AI가 아니라 실제로 파일을 열어본 사람이 찾아냈다", 1),
    "촬영 도구가 지정 해상도를 무시 → 헤드리스 Chrome CLI에서 Playwright로 교체",
    ("Playwright 도입으로 첫화면/진행중/종료 3상태 촬영까지 함께 가능해짐", 1),
    "제출용 PDF에서 이모지가 네모 상자로 깨짐 → 버튼 라벨에서 이모지 제거",
    ("'내 환경에선 정상'은 제출 증거에서는 판단 기준이 될 수 없다는 교훈", 1),
], size=17)

# ═══════════════════════════════════════════════════════════
# 17. 최종 체크리스트
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "FINAL STATUS", "제출 준비 완료")
rows = [
    ["카드 1~5 통과 기준", "5/5", "완료"],
    ["학생 자체 점검", "11/11", "완료"],
    ["콘솔 오류", "0건", "완료"],
    ["개인정보·비밀값", "0건 (실질)", "완료"],
    ["10분 연속 실행", "DOM·힙 증가 0", "완료"],
    ["제출물(검증안내서·AI3줄)", "작성 완료", "완료"],
]
table_slide(s, ["항목", "결과", "상태"], rows,
            Inches(2.4), Inches(1.9), Inches(8.5), Inches(4.2),
            col_w=[Inches(4.5), Inches(2.7), Inches(1.3)], font_size=16)

# ═══════════════════════════════════════════════════════════
# 17d. 검증 안내서 — 개요 (주소 + 3단계)
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "VERIFICATION GUIDE", "검증 안내서 — 직접 확인하는 방법")
_, tf = box(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.9))
para(tf, "배포 주소", 14, color=INK_DIM, first=True, space_after=2)
para(tf, URL, 24, color=BRAND, bold=True)
bullets(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(3.4), [
    "1단계 — 위 주소로 접속해 [▶ 시작하기]를 누른다",
    "2단계 — 격자에서 딱 하나 다른 칸을 클릭(터치)한다",
    ("키보드로 하려면 방향키로 이동한 뒤 Enter", 1),
    "3단계 — 시간이 0이 되어 결과 화면이 뜨면 [다시 하기]를 누른다",
], size=21, gap=18)
_, tf2 = box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6))
para(tf2, "3단계 · 30초 이내 — 설치·로그인 없이 링크만으로 확인 가능", 15, color=INK_DIM, first=True)

# ═══════════════════════════════════════════════════════════
# 17e. 검증 안내서 — 통과 기준 (1~5)
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "VERIFICATION GUIDE", "통과 기준 — 이렇게 보이면 정상입니다 (1~5)")
rows = [
    ["규칙·조작·상태 상시 표시", "화면 하단 3줄이 항상 보임 — 규칙 · 조작 · 현재 상태(단계·난이도), 진행에 따라 갱신"],
    ["조작 1회 = 결과 1회", "정답 → 즉시 다음 격자 + STAGE·SCORE 상승 · 오답 → 칸이 흔들리고 시간 1초 차감(게임 계속) · 시간 초과 → 정답 위치가 민트 테두리로 공개"],
    ["다시 시작 = 초기화", "[다시 하기] 시 STAGE 1 · SCORE 0으로 복귀, 단 최고 기록은 그대로 보존"],
    ["음소거 작동", "⚙ → 소리 끄면 즉시 무음, 다시 켜면 즉시 복원"],
    ["움직임 줄이기 작동", "⚙ → 켜면 확대·흔들림 대신 색 점멸로 대체 (성공/실패 인지는 그대로 유지)"],
]
table_slide(s, ["확인 항목", "이렇게 보이면 통과"], rows,
            Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.9),
            col_w=[Inches(3.2), Inches(8.7)], font_size=14, left_cols={0, 1})

# ═══════════════════════════════════════════════════════════
# 17f. 검증 안내서 — 통과 기준 (6~9) + 추가로 볼 수 있는 것
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "VERIFICATION GUIDE", "통과 기준 — 이렇게 보이면 정상입니다 (6~9)")
rows = [
    ["일시정지 작동", "P·Esc·❙❙ 클릭 시 타이머 정지, 정지 중 칸 클릭 무반응 · [계속하기] 시 남은 시간 그대로 이어짐"],
    ["같은 링크 = 같은 문제", "[결과 공유]로 나온 ?c=코드 링크를 다른 기기·브라우저에서 열어도 1단계부터 완전히 동일"],
    ["콘솔 오류 없음", "F12 → Console 탭에 빨간 오류 0건"],
    ["두 해상도에서 안 잘림", "1366×768 / 1920×1080 모두 가로 스크롤 없이 격자·하단 3줄이 화면 안에 들어옴"],
]
table_slide(s, ["확인 항목", "이렇게 보이면 통과"], rows,
            Inches(0.7), Inches(1.85), Inches(11.9), Inches(2.5),
            col_w=[Inches(3.2), Inches(8.7)], font_size=14, left_cols={0, 1})
_, tf = box(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.35))
para(tf, "추가로 볼 수 있는 것 (선택)", 15, color=BRAND, bold=True, first=True)
bullets(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(2.0), [
    "오늘의 문제 — [오늘의 문제] 클릭 시 그날 날짜(KST 기준)로 고정된 문제, 같은 날이면 누구나 동일",
    "결과 카드 — 점수 + 마지막 문제 격자를 정답 표시 없이 포함해 보는 사람이 직접 풀어볼 수 있음",
    "차이의 종류 — 정답 시 화면 가운데에 `회전 · 0.42초`처럼 어떤 차이였는지 표시 (10종)",
], size=15)

# ═══════════════════════════════════════════════════════════
# 17g. 검증 안내서 — 안 될 때 (FAQ)
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
title_bar(s, "VERIFICATION GUIDE", "안 될 때 — 자주 묻는 질문")
rows = [
    ["화면이 비어 있다", "Ctrl+F5로 강력 새로고침"],
    ["로컬에서 열었더니 안 나온다", "index.html 더블클릭은 동작 안 함 — ES 모듈은 file://에서 차단되므로 서버로 열 것"],
    ["소리가 안 난다", "브라우저는 사용자 입력 전엔 소리를 못 냄 — [시작하기] 클릭 후에도 안 나면 ⚙ 소리 설정·시스템 볼륨 확인"],
    ["키보드를 눌러도 반응이 없다", "포커스가 버튼에 남아있을 수 있음 — 격자 칸을 한 번 클릭한 뒤 방향키 사용"],
    ["게임이 저절로 멈춘 것 같다", "P나 Esc를 눌렀는지 확인 — 일시정지 화면에서 [계속하기] 클릭"],
    ["기록이 이상하다", "⚙ → [기록 초기화] — 저장값이 깨져도 게임은 기본값으로 정상 실행됨"],
]
table_slide(s, ["증상", "해결 방법"], rows,
            Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.9),
            col_w=[Inches(3.4), Inches(8.5)], font_size=13.5, left_cols={0, 1})

# ═══════════════════════════════════════════════════════════
# 17b/17c. 필수 체크리스트 — 11항목, 스크린샷 증거 갤러리 (2장)
# ═══════════════════════════════════════════════════════════
D = "img/difficulty"; A = "img/actions"; ST = "img/stability"; EV = "img/evidence"

items_p1 = [
    ("1", "시작→종료→재시작", [(f"{A}/01_start_after.png","시작"), (f"{A}/08_gameover_after.png","종료")], None),
    ("2", "성공·실패 상태 구분", [(f"{A}/02_correct_after.png","정답→다음단계"), (f"{A}/03_wrong_after.png","오답→시간차감")], None),
    ("3", "2해상도 잘림·넘침 없음", [(f"{EV}/resolution_1366x768.png","1366×768"), (f"{EV}/resolution_1920x1080.png","1920×1080")], None),
    ("4", "연속입력·포커스·10분", [(f"{A}/04_pause_after.png","일시정지"), (f"{ST}/t600_10min.png","10분 도달")], None),
    ("5", "콘솔 빨간 오류 0건", [("img/콘솔오류확인.png","실기기 DevTools — 배포 주소")], "사용자 직접 확인 · 최고 16단계·3,242점"),
    ("6", "난이도 전/후 10+10", [(f"{D}/result_T4_4.png","T=4.4 결과"), (f"{D}/result_T5_4.png","T=5.4 결과")], None),
]
items_p2 = [
    ("7", "손상값 → 기본값 실행", [(f"{EV}/storage_before_normal.png","손상 전(최고 5단계)"), (f"{EV}/storage_after_corrupted_still_works.png","손상 후(정상 실행)")], None),
    ("8", "재시작=초기화, 기록 유지", [(f"{A}/13_restart_before.png","결과 화면"), (f"{A}/13_restart_after.png","초기화+최고기록")], None),
    ("9", "효과 줄이기·끄기 작동", [(f"{A}/06_sound_after.png","소리 끄기"), (f"{A}/07_motion_after.png","움직임 줄이기")], None),
    ("10", "개인정보·비밀값 0건", [], "저장소 전수 grep 검사 — stulss(공개 계정명) 외 유출 0건"),
    ("11", "위 1~10 모두 완료", [], "카드 1~5 통과기준 5/5 · 학생 자체 점검 11/11"),
]

def checklist_gallery_slide(kicker, title, items):
    s = track(new_slide())
    title_bar(s, kicker, title)
    cols, rows_n = 3, 2
    gx0, gy0 = Inches(0.5), Inches(1.75)
    cw, ch, ggap = Inches(4.0), Inches(2.62), Inches(0.15)
    for i, (num, t, imgs, note) in enumerate(items):
        r, c = divmod(i, cols)
        x = gx0 + c * (cw + ggap); y = gy0 + r * (ch + ggap)
        evidence_cell(s, x, y, cw, ch, num, t, imgs, note)
    return s

checklist_gallery_slide("필수 체크리스트 · 증거 1/2", "학생 자체 점검 — 스크린샷 증거 (1~6)", items_p1)
checklist_gallery_slide("필수 체크리스트 · 증거 2/2", "학생 자체 점검 — 스크린샷 증거 (7~11)", items_p2)

# ═══════════════════════════════════════════════════════════
# 18. 마무리
# ═══════════════════════════════════════════════════════════
s = track(new_slide())
_, tf = box(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.0))
para(tf, "감사합니다", 44, color=INK, bold=True, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6))
para(tf, "▶ 플레이: " + URL, 17, color=BRAND, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.5))
para(tf, "소스: " + REPO, 13, color=INK_DIM, align=PP_ALIGN.CENTER, first=True)

OUT = "발표자료.pptx"
prs.save(OUT)
print(f"생성 완료: {OUT} ({os.path.getsize(OUT)/1024:.0f}KB) · 슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장" if False else
      f"생성 완료: {OUT} ({os.path.getsize(OUT)/1024:.0f}KB) · 슬라이드 {N[0]}장")
